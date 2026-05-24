"""
App: AI Estudos+ — Gerador de Artigos, Planos de Aula e Slides

Transforma o exemplo simples em uma aplicação completa e útil para estudo, docência
(e produção acadêmica). Mantém streaming via Groq e adiciona:

1) Aba Artigo Acadêmico (com esqueleto/outline, reescrita de seções e exportação DOCX/MD)
2) Aba Plano de Aula (estrutura didática guiada por objetivos, turma e duração)
3) Aba Slides (gera um .pptx a partir de um resumo)
4) Aba Configurações (modelos, temperatura, chaves)

Requisitos (requirements.txt):
- gradio
- groq
- python-dotenv
- python-docx
- python-pptx

Variáveis de ambiente esperadas (ou informadas na UI em "Configurações"):
- GROQ_API_KEY=<sua_chave>

Como executar:
$ python app_ai_estudos_plus.py
"""
from __future__ import annotations

import os
import io
import re
from datetime import datetime
from typing import Generator, Optional, Tuple

import gradio as gr
from groq import Groq
from dotenv import load_dotenv
from docx import Document
from docx.shared import Pt
from pptx import Presentation

# -------------------------
# Config + Helpers
# -------------------------
load_dotenv()

DEFAULT_MODEL = "openai/gpt-oss-120b"  # mantém compatível com o exemplo original
SUPPORTED_MODELS = [
    "openai/gpt-oss-120b",
]


def get_client(api_key: Optional[str] = None) -> Groq:
    key = api_key or os.getenv("GROQ_API_KEY")
    if not key:
        raise RuntimeError(
            "Defina a GROQ_API_KEY no .env ou informe na aba Configurações.")
    return Groq(api_key=key)


def chat_stream(
    prompt: str,
    model: str,
    temperature: float = 0.8,
    max_tokens: int = 4000,
    system: Optional[str] = None,
    api_key: Optional[str] = None,
) -> Generator[str, None, None]:
    """Stream de texto usando Groq Chat Completions."""
    client = get_client(api_key)
    messages = []
    if system:
        messages.append({"role": "system", "content": system})
    messages.append({"role": "user", "content": prompt})

    completion = client.chat.completions.create(
        model=model,
        messages=messages,
        temperature=float(temperature),
        max_completion_tokens=int(max_tokens),
        top_p=1,
        reasoning_effort="medium",
        stream=True,
        stop=None,
    )

    acc = ""
    for chunk in completion:
        delta = chunk.choices[0].delta.content or ""
        acc += delta
        yield acc


# -------------------------
# Geração de conteúdo (prompts guiados)
# -------------------------

SYSTEM_ACADEMIC = (
    "Você é um assistente acadêmico que escreve de forma clara, objetiva e sem plágio, "
    "com linguagem adequada ao nível indicado, usando exemplos e seções bem organizadas.\n"
)

SYSTEM_LESSON = (
    "Você é um designer instrucional. Produza planos de aula acionáveis, com objetivos, "
    "etapas e instrumentos de avaliação, adequados ao nível indicado.\n"
)

SYSTEM_SLIDES = (
    "Você transforma resumos em tópicos para slides concisos (título + 3-6 bullets), "
    "com mensagens claras e linguagem simples.\n"
)


def build_outline_prompt(
    tema: str,
    palavras_chave: str,
    nivel: str,
    formato: str,
    tamanho_palavras: int,
) -> str:
    return f"""
Crie um esqueleto (outline) de artigo acadêmico em {formato} sobre o tema: "{tema}".
Inclua títulos de seções e subtítulos com bullets do conteúdo esperado.
Nível do leitor: {nivel}. Tamanho alvo: ~{tamanho_palavras} palavras.
Use estrutura típica: Resumo, Introdução, Fundamentação/Referencial, Metodologia (se aplicável), Resultados/Discussão,
Conclusão e Referências (com placeholders [REF] — não invente dados específicos).
Palavras-chave: {palavras_chave}
""".strip()


def build_article_prompt(
    outline: str,
    objetivo: str,
    formato: str,
    nivel: str,
    tamanho_palavras: int,
    estilo: str,
) -> str:
    return f"""
Com base no outline abaixo, escreva o artigo completo em {formato}.
- Objetivo do texto: {objetivo}
- Nível do leitor: {nivel}
- Tamanho aproximado: {tamanho_palavras} palavras
- Estilo: {estilo}

Outline:
{outline}

Diretrizes:
- Escreva com clareza e coesão; use seções e subtítulos.
- Não invente citações reais; coloque placeholders tipo [Autor, Ano] quando necessário.
- Inclua uma seção final "Referências" apenas com placeholders (não gere DOIs ou URLs falsos).
""".strip()


def build_rewrite_prompt(
    trecho: str, instrucao: str, nivel: str, formato: str
) -> str:
    return f"""
Reescreva o trecho abaixo aplicando as instruções, mantendo o formato {formato} e
linguagem para nível {nivel}. Melhore clareza, concisão e fluidez. Preserve o sentido.

Instruções do autor: {instrucao}

Trecho:
{trecho}
""".strip()


def build_lesson_plan_prompt(
    tema: str,
    serie: str,
    duracao_min: int,
    objetivo: str,
    metodologia: str,
    avaliacao: str,
) -> str:
    return f"""
Crie um plano de aula completo sobre "{tema}" para a turma/série "{serie}".
Duração: {duracao_min} minutos. Objetivo: {objetivo}.
Inclua: Contextualização, Objetivos de Aprendizagem (em verbos observáveis), Materiais,
Etapas com tempo (Abertura, Desenvolvimento, Consolidação), Atividades práticas,
Diferenciação (estratégias para alunos com diferentes níveis) e Avaliação ({avaliacao}).
Metodologia sugerida: {metodologia}.
No final, inclua uma seção "Tarefa/Extensão" e "Rubrica de Avaliação" com critérios e níveis.
""".strip()


def build_slides_prompt(resumo: str, publico: str, num_slides: int) -> str:
    return f"""
Transforme o resumo abaixo em tópicos para aproximadamente {num_slides} slides, voltados ao público "{publico}".
Cada slide deve conter: Título curto + 3 a 6 bullets concisos.

Resumo:
{resumo}
""".strip()


# -------------------------
# Exporters
# -------------------------

def sanitize_filename(name: str) -> str:
    name = re.sub(r"[\\/:*?\"<>|]", "_", name).strip()
    return name or "documento"


def export_docx(title: str, text: str) -> bytes:
    doc = Document()
    # Título simples
    h = doc.add_heading(title.strip() or "Documento", level=0)
    for run in h.runs:
        run.font.size = Pt(18)
    # Corpo
    for line in text.splitlines():
        doc.add_paragraph(line)
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


def export_md(title: str, text: str) -> bytes:
    md = f"# {title.strip() or 'Documento'}\n\n{text}\n"
    return md.encode("utf-8")


def export_pptx(title: str, slides_text: str) -> bytes:
    prs = Presentation()
    # Título
    layout_title = prs.slide_layouts[0]
    slide0 = prs.slides.add_slide(layout_title)
    slide0.shapes.title.text = title.strip() or "Apresentação"
    slide0.placeholders[1].text = datetime.now().strftime("Gerado em %d/%m/%Y")

    # Demais slides (espera-se um formato "\n\n# Slide Title\n- bullet\n- bullet")
    blocks = re.split(r"\n\s*\n", slides_text.strip())
    for block in blocks:
        lines = [l for l in block.splitlines() if l.strip()]
        if not lines:
            continue
        stitle = lines[0].lstrip("# ").strip() if lines[0].startswith("#") else lines[0][:60]
        content_bullets = [l.lstrip("- ").strip() for l in lines[1:]]

        layout = prs.slide_layouts[1]
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = stitle
        body = s.placeholders[1].text_frame
        body.clear()
        for i, b in enumerate(content_bullets):
            p = body.add_paragraph() if i > 0 else body.paragraphs[0]
            p.text = b
            p.level = 0

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# -------------------------
# UI — Gradio
# -------------------------
with gr.Blocks(title="AI Estudos+") as demo:
    gr.Markdown("""
    # 📚 AI Estudos+
    Um estúdio de escrita acadêmica e ensino com **Groq** (streaming), pronto para gerar:
    - **Artigos acadêmicos** com outline, reescrita de seções e exportações (.docx/.md)
    - **Planos de aula** completos e acionáveis
    - **Slides** (.pptx) a partir de um resumo estruturado
    """)

    # Estados globais
    st_outline = gr.State("")
    st_model = gr.State(DEFAULT_MODEL)
    st_api_key = gr.State(os.getenv("GROQ_API_KEY", ""))

    with gr.Tab("Artigo Acadêmico"):
        with gr.Row():
            with gr.Column(scale=1):
                tema = gr.Textbox(label="Tema", placeholder="Ex.: Impactos da IA na Educação Básica")
                palavras = gr.Textbox(label="Palavras-chave", placeholder="IA, educação, avaliação formativa")
                objetivo = gr.Textbox(label="Objetivo do texto", value="Explicar conceitos e discutir implicações práticas")
                nivel = gr.Dropdown(["Fundamental", "Médio", "Graduação", "Pós"], value="Graduação", label="Nível do leitor")
                formato = gr.Dropdown(["ABNT", "APA", "MLA"], value="ABNT", label="Formato")
                tamanho = gr.Slider(600, 5000, value=1200, step=100, label="Tamanho alvo (palavras)")
                estilo = gr.Dropdown(["Clássico", "Objetivo", "Didático", "Crítico"], value="Didático", label="Estilo")
                btn_outline = gr.Button("🧱 Gerar Outline", variant="secondary")
                btn_article = gr.Button("✍️ Gerar Artigo", variant="primary")
                with gr.Accordion("Reescrever trecho", open=False):
                    trecho = gr.Textbox(label="Cole um trecho para melhorar", lines=6)
                    instrucao = gr.Textbox(label="Como quer melhorar?", value="Mais claro, menos jargão, exemplos práticos.")
                    btn_rewrite = gr.Button("🔁 Reescrever trecho")

                with gr.Row():
                    title_export = gr.Textbox(label="Título para exportação", value="Artigo")
                with gr.Row():
                    btn_export_docx = gr.Button("⬇️ Baixar DOCX")
                    btn_export_md = gr.Button("⬇️ Baixar Markdown")

            with gr.Column(scale=2):
                gr.Markdown("### Resultado / Edição")
                editor = gr.Textbox(lines=28, label="Editor de texto", placeholder="O conteúdo aparecerá aqui e pode ser editado.")
                word_count = gr.Markdown("**Palavras:** 0")

        # Callbacks — Outline
        def do_outline(tema, palavras, nivel, formato, tamanho, model, api_key):
            prompt = build_outline_prompt(tema, palavras, nivel, formato, int(tamanho))
            text = ""
            for partial in chat_stream(prompt, model=model, temperature=0.7, max_tokens=1500, system=SYSTEM_ACADEMIC, api_key=api_key):
                text = partial
                yield text, text  # editor e state

        btn_outline.click(
            fn=do_outline,
            inputs=[tema, palavras, nivel, formato, tamanho, st_model, st_api_key],
            outputs=[editor, st_outline],
        )

        # Callbacks — Artigo
        def do_article(outline, objetivo, formato, nivel, tamanho, estilo, model, api_key):
            if not outline.strip():
                outline = "(Sem outline — gere primeiro ou cole um)"
            prompt = build_article_prompt(outline, objetivo, formato, nivel, int(tamanho), estilo)
            text = ""
            for partial in chat_stream(prompt, model=model, temperature=0.85, max_tokens=3000, system=SYSTEM_ACADEMIC, api_key=api_key):
                text = partial
                yield text

        btn_article.click(
            fn=do_article,
            inputs=[st_outline, objetivo, formato, nivel, tamanho, estilo, st_model, st_api_key],
            outputs=editor,
        )

        # Reescrever trecho
        def do_rewrite(trecho, instrucao, nivel, formato, model, api_key):
            if not trecho.strip():
                return "Cole um trecho para reescrever."
            prompt = build_rewrite_prompt(trecho, instrucao, nivel, formato)
            text = ""
            for partial in chat_stream(prompt, model=model, temperature=0.7, max_tokens=1200, system=SYSTEM_ACADEMIC, api_key=api_key):
                text = partial
                yield text

        btn_rewrite.click(
            fn=do_rewrite,
            inputs=[trecho, instrucao, nivel, formato, st_model, st_api_key],
            outputs=editor,
        )

        # Contador de palavras
        def count_words(txt: str) -> str:
            n = len([w for w in re.findall(r"\b\w+\b", txt or "")])
            return f"**Palavras:** {n}"

        editor.change(fn=count_words, inputs=editor, outputs=word_count)

        # Exportações — agora salvando no disco e retornando o caminho
        def do_export_docx(title: str, text: str) -> str:
            filename = sanitize_filename(title) + ".docx"
            with open(filename, "wb") as f:
                f.write(export_docx(title, text))
            return filename

        def do_export_md(title: str, text: str) -> str:
            filename = sanitize_filename(title) + ".md"
            with open(filename, "wb") as f:
                f.write(export_md(title, text))
            return filename

        out_docx = gr.File(label="Documento .docx")
        out_md = gr.File(label="Documento .md")
        btn_export_docx.click(fn=do_export_docx, inputs=[title_export, editor], outputs=out_docx)
        btn_export_md.click(fn=do_export_md, inputs=[title_export, editor], outputs=out_md)

    # ------------------
    # Plano de Aula
    # ------------------
    with gr.Tab("Plano de Aula"):
        with gr.Row():
            with gr.Column(scale=1):
                tema_l = gr.Textbox(label="Tema da aula", value="Regressão Linear (introdução)")
                serie_l = gr.Textbox(label="Série/Turma", value="Ensino Médio — 2º ano")
                duracao_l = gr.Slider(20, 200, value=50, step=5, label="Duração (min)")
                objetivo_l = gr.Textbox(label="Objetivo de aprendizagem", value="Compreender a ideia de ajustar uma reta a dados")
                metodologia_l = gr.Textbox(label="Metodologia", value="Exposição dialogada + atividade prática em pares")
                avaliacao_l = gr.Textbox(label="Avaliação", value="Checklist de participação e mini-quiz de 5 questões")
                btn_plan = gr.Button("📋 Gerar Plano")
            with gr.Column(scale=2):
                plan_out = gr.Textbox(lines=26, label="Plano de aula (editável)")
                with gr.Row():
                    title_plan = gr.Textbox(label="Título do arquivo", value="plano_de_aula")
                    btn_plan_docx = gr.Button("⬇️ Baixar DOCX")
                plan_docx = gr.File(label="Plano .docx")

        def do_plan(tema, serie, duracao, objetivo, metodologia, avaliacao, model, api_key):
            prompt = build_lesson_plan_prompt(tema, serie, int(duracao), objetivo, metodologia, avaliacao)
            text = ""
            for partial in chat_stream(prompt, model=model, temperature=0.7, max_tokens=2000, system=SYSTEM_LESSON, api_key=api_key):
                text = partial
                yield text

        btn_plan.click(
            fn=do_plan,
            inputs=[tema_l, serie_l, duracao_l, objetivo_l, metodologia_l, avaliacao_l, st_model, st_api_key],
            outputs=plan_out,
        )

        def export_plan_docx(title: str, text: str) -> str:
            filename = sanitize_filename(title) + ".docx"
            with open(filename, "wb") as f:
                f.write(export_docx(title, text))
            return filename

        btn_plan_docx.click(export_plan_docx, inputs=[title_plan, plan_out], outputs=plan_docx)

    # ------------------
    # Slides
    # ------------------
    with gr.Tab("Slides (.pptx)"):
        with gr.Row():
            with gr.Column(scale=1):
                titulo_s = gr.Textbox(label="Título da apresentação", value="Introdução à Regressão Linear")
                publico_s = gr.Textbox(label="Público-alvo", value="Iniciantes em Ciência de Dados")
                resumo_s = gr.Textbox(label="Resumo do conteúdo", lines=12, value=(
                    "Definição de regressão linear, equação da reta, objetivo de minimizar erro.\n"
                    "Exemplos simples, interpretação de coeficientes, overfitting vs underfitting.\n"
                    "Demonstração rápida com biblioteca.")
                )
                num_slides = gr.Slider(4, 25, value=8, step=1, label="Nº aproximado de slides")
                btn_slides_outline = gr.Button("🗂️ Rascunhar Tópicos")
                btn_build_pptx = gr.Button("🎞️ Gerar .pptx")
            with gr.Column(scale=2):
                slides_txt = gr.Textbox(lines=22, label="Tópicos para slides (edite antes de gerar o .pptx)")
                pptx_file = gr.File(label="Slides .pptx")

        def do_slides_outline(resumo, publico, n, model, api_key):
            prompt = build_slides_prompt(resumo, publico, int(n))
            text = ""
            for partial in chat_stream(prompt, model=model, temperature=0.6, max_tokens=1200, system=SYSTEM_SLIDES, api_key=api_key):
                text = partial
                yield text

        btn_slides_outline.click(
            fn=do_slides_outline,
            inputs=[resumo_s, publico_s, num_slides, st_model, st_api_key],
            outputs=slides_txt,
        )

        def do_build_pptx(title: str, slides: str) -> str:
            filename = sanitize_filename(title) + ".pptx"
            with open(filename, "wb") as f:
                f.write(export_pptx(title, slides))
            return filename

        btn_build_pptx.click(fn=do_build_pptx, inputs=[titulo_s, slides_txt], outputs=pptx_file)

    # ------------------
    # Configurações
    # ------------------
    with gr.Tab("Configurações"):
        gr.Markdown("Forneça sua chave (opcional se já estiver no .env) e escolha o modelo.")
        api_key_in = gr.Textbox(type="password", label="GROQ_API_KEY", value=os.getenv("GROQ_API_KEY", ""))
        model_dd = gr.Dropdown(SUPPORTED_MODELS, value=DEFAULT_MODEL, label="Modelo")
        temp_def = gr.Slider(0.0, 2.0, value=0.8, step=0.05, label="Temperatura padrão (usada pelos prompts)")
        btn_apply = gr.Button("Aplicar")

        def apply_settings(key, model, temp):
            gr.Info(f"Configurações aplicadas. Modelo: {model} | Temp: {temp}")
            return key, model

        btn_apply.click(fn=apply_settings, inputs=[api_key_in, model_dd, temp_def], outputs=[st_api_key, st_model])


if __name__ == "__main__":
    demo.launch()  # use share=True se quiser link público no Gradio
